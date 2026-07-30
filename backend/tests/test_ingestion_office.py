"""DOCX, PPTX, XLSX, CSV and HTML parsing.

Every fixture is a genuine file produced by the same library that reads it
back, so these exercise the real extraction path rather than a canned string.
"""

import csv
import io
import zipfile

import docx
import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.ingestion import get_parser
from app.ingestion.base import UnparsableDocumentError
from app.ingestion.powerpoint import PptxParser
from app.ingestion.tabular import CsvParser, XlsxParser
from app.ingestion.web import HtmlParser
from app.ingestion.word import DocxParser

# ==========================================================================
# Builders
# ==========================================================================


def build_docx(*, with_page_break: bool = False, with_table: bool = False) -> bytes:
    document = docx.Document()
    document.add_heading("Quarterly Report", level=1)
    document.add_paragraph("Revenue grew by twelve percent this quarter.")

    if with_table:
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Region"
        table.cell(0, 1).text = "Revenue"
        table.cell(1, 0).text = "North"
        table.cell(1, 1).text = "44000"

    if with_page_break:
        document.add_page_break()
        document.add_paragraph("Appendix content lives after the break.")

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_pptx(*, with_notes: bool = False) -> bytes:
    presentation = Presentation()

    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Architecture Overview"
    first.placeholders[1].text = "Retrieval augmented generation pipeline"

    if with_notes:
        first.notes_slide.notes_text_frame.text = "Mention the reranking step here."

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Results"
    box = second.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Latency dropped to 45 milliseconds"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_xlsx() -> bytes:
    workbook = openpyxl.Workbook()

    sheet = workbook.active
    sheet.title = "Revenue"
    sheet.append(["Region", "Quarter", "Amount"])
    sheet.append(["North", "Q3", 44000])
    sheet.append(["South", "Q3", 31000])

    second = workbook.create_sheet("Headcount")
    second.append(["Team", "People"])
    second.append(["Platform", 12])

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_csv(delimiter: str = ",") -> bytes:
    buffer = io.StringIO()
    writer = csv.writer(buffer, delimiter=delimiter)
    writer.writerow(["Product", "SKU", "Price"])
    writer.writerow(["Keyboard", "4471", "89.00"])
    writer.writerow(["Monitor", "9920", "310.50"])
    return buffer.getvalue().encode()


HTML_PAGE = b"""<!doctype html>
<html><head><title>Indexing Guide</title>
<style>body { color: red; }</style>
<script>console.log('tracking');</script></head>
<body>
  <nav><a href="/">Home</a><a href="/docs">Docs</a></nav>
  <main>
    <h1>B-tree indexes</h1>
    <p>A B-tree index speeds up equality and range queries.</p>
    <ul><li>Use partial indexes for repeated predicates</li></ul>
  </main>
  <footer>Copyright 2026</footer>
</body></html>"""


# ==========================================================================
# DOCX
# ==========================================================================


def test_docx_text_is_extracted() -> None:
    parsed = DocxParser().parse(build_docx(), filename="report.docx")

    assert parsed.page_count == 1
    assert "Revenue grew by twelve percent" in parsed.pages[0].text


def test_docx_headings_become_markdown() -> None:
    """Structure must survive so the chunker can split on section boundaries."""
    parsed = DocxParser().parse(build_docx(), filename="report.docx")

    assert "# Quarterly Report" in parsed.pages[0].text


def test_docx_explicit_page_breaks_create_pages() -> None:
    """Word stores no pagination, so an author's page break is the only
    honest boundary available for a citation."""
    parsed = DocxParser().parse(build_docx(with_page_break=True), filename="report.docx")

    assert parsed.page_count == 2
    assert "Appendix" in parsed.pages[1].text
    assert [page.number for page in parsed.pages] == [1, 2]


def test_docx_tables_keep_column_names_with_values() -> None:
    """A bare row loses the association between "44000" and "Revenue"."""
    parsed = DocxParser().parse(build_docx(with_table=True), filename="report.docx")

    assert "Region: North" in parsed.pages[0].text
    assert "Revenue: 44000" in parsed.pages[0].text


def test_docx_rejects_a_non_docx_file() -> None:
    """The container check runs before python-docx, so a file that only
    pretends to be a ZIP is caught there with a clearer message."""
    with pytest.raises(UnparsableDocumentError, match="container is corrupt"):
        DocxParser().parse(b"PK\x03\x04 not really a docx", filename="fake.docx")


def test_docx_rejects_a_valid_zip_that_is_not_a_word_file() -> None:
    """Past the container check, python-docx's own rejection must surface —
    with guidance, since .doc is the common cause."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("hello.txt", b"this is a zip, but not a docx")

    with pytest.raises(UnparsableDocumentError, match="Word document"):
        DocxParser().parse(buffer.getvalue(), filename="fake.docx")


# ==========================================================================
# PPTX
# ==========================================================================


def test_pptx_slides_map_to_pages() -> None:
    parsed = PptxParser().parse(build_pptx(), filename="deck.pptx")

    assert parsed.page_count == 2
    assert [page.number for page in parsed.pages] == [1, 2]
    assert "Architecture Overview" in parsed.pages[0].text
    assert "Latency dropped" in parsed.pages[1].text


def test_pptx_slide_titles_become_headings() -> None:
    parsed = PptxParser().parse(build_pptx(), filename="deck.pptx")

    assert parsed.pages[0].text.startswith("# Architecture Overview")


def test_pptx_speaker_notes_are_captured() -> None:
    """Slides carry fragments; the notes often carry the actual argument."""
    parsed = PptxParser().parse(build_pptx(with_notes=True), filename="deck.pptx")

    assert "Speaker notes: Mention the reranking step" in parsed.pages[0].text


# ==========================================================================
# XLSX
# ==========================================================================


def test_xlsx_sheets_map_to_pages() -> None:
    parsed = XlsxParser().parse(build_xlsx(), filename="figures.xlsx")

    assert parsed.page_count == 2
    assert "Sheet: Revenue" in parsed.pages[0].text
    assert "Sheet: Headcount" in parsed.pages[1].text


def test_xlsx_rows_are_rendered_with_headers() -> None:
    parsed = XlsxParser().parse(build_xlsx(), filename="figures.xlsx")

    assert "Region: North | Quarter: Q3 | Amount: 44000" in parsed.pages[0].text


def test_xlsx_integers_do_not_become_floats() -> None:
    """openpyxl returns numbers as floats, so a year would embed as "2024.0"."""
    parsed = XlsxParser().parse(build_xlsx(), filename="figures.xlsx")

    assert "44000.0" not in parsed.pages[0].text
    assert "44000" in parsed.pages[0].text


def test_xlsx_rejects_garbage() -> None:
    with pytest.raises(UnparsableDocumentError):
        XlsxParser().parse(b"definitely not a spreadsheet", filename="fake.xlsx")


# ==========================================================================
# CSV
# ==========================================================================


def test_csv_rows_are_rendered_with_headers() -> None:
    parsed = CsvParser().parse(build_csv(), filename="products.csv")

    assert "Product: Keyboard | SKU: 4471 | Price: 89.00" in parsed.pages[0].text


def test_csv_semicolon_delimiter_is_detected() -> None:
    """European exports use semicolons; assuming commas yields one column."""
    parsed = CsvParser().parse(build_csv(delimiter=";"), filename="products.csv")

    assert "Product: Keyboard" in parsed.pages[0].text


def test_csv_is_a_single_page() -> None:
    parsed = CsvParser().parse(build_csv(), filename="products.csv")

    assert parsed.page_count == 1
    assert parsed.pages[0].number == 1


def test_empty_csv_is_rejected() -> None:
    with pytest.raises(UnparsableDocumentError, match="empty"):
        CsvParser().parse(b"", filename="empty.csv")


# ==========================================================================
# HTML
# ==========================================================================


def test_html_content_is_extracted() -> None:
    parsed = HtmlParser().parse(HTML_PAGE, filename="guide.html")

    assert "B-tree index speeds up equality" in parsed.pages[0].text


def test_html_scripts_styles_and_chrome_are_stripped() -> None:
    """Menu labels and cookie notices match every query weakly and crowd out
    real passages — the same failure mode as the bibliography problem."""
    text = HtmlParser().parse(HTML_PAGE, filename="guide.html").pages[0].text

    assert "console.log" not in text
    assert "color: red" not in text
    assert "Copyright 2026" not in text
    assert "Docs" not in text


def test_html_title_and_headings_are_preserved() -> None:
    text = HtmlParser().parse(HTML_PAGE, filename="guide.html").pages[0].text

    assert text.startswith("# Indexing Guide")
    assert "# B-tree indexes" in text


def test_html_list_items_are_bulleted() -> None:
    text = HtmlParser().parse(HTML_PAGE, filename="guide.html").pages[0].text

    assert "- Use partial indexes" in text


def test_html_with_no_text_is_rejected() -> None:
    with pytest.raises(UnparsableDocumentError, match="JavaScript"):
        HtmlParser().parse(b"<html><body><script>x=1</script></body></html>", filename="app.html")


# ==========================================================================
# Zip-bomb guard
# ==========================================================================


def test_decompression_bomb_is_rejected() -> None:
    """A small archive that expands enormously must not reach the parser.

    The upload limit bounds what arrives; this bounds what it becomes.
    """
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", b"\0" * (60 * 1024 * 1024))

    with pytest.raises(UnparsableDocumentError, match="compression ratio|400 MB"):
        DocxParser().parse(buffer.getvalue(), filename="bomb.docx")


# ==========================================================================
# Registry
# ==========================================================================


@pytest.mark.parametrize(
    ("filename", "expected"),
    [
        ("a.docx", DocxParser),
        ("a.pptx", PptxParser),
        ("a.xlsx", XlsxParser),
        ("a.xlsm", XlsxParser),
        ("a.csv", CsvParser),
        ("a.tsv", CsvParser),
        ("a.html", HtmlParser),
        ("a.htm", HtmlParser),
    ],
)
def test_registry_routes_by_extension(filename: str, expected: type) -> None:
    assert isinstance(get_parser(filename=filename, content_type=None), expected)


def test_csv_is_not_captured_by_the_text_parser() -> None:
    """TextParser claims text/plain, which browsers report for CSV files.

    Extension matching runs first precisely so a .csv still gets row-aware
    rendering rather than being dumped as raw lines.
    """
    parser = get_parser(filename="products.csv", content_type="text/plain")

    assert isinstance(parser, CsvParser)


def test_legacy_office_formats_are_rejected_with_guidance() -> None:
    with pytest.raises(UnparsableDocumentError, match="Unsupported"):
        get_parser(filename="old.doc", content_type=None)
