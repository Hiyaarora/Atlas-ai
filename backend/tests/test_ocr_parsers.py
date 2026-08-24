"""OCR inside the real parsers: PDF, DOCX and PPTX.

Documents are built in memory rather than committed as fixtures, so each test
states exactly what is in its file and cannot drift from a binary nobody
opens.

The OCR engine is stubbed. These verify the *wiring* — that images are found,
that recognised text lands on the right page carrying the right metadata, and
above all that existing text extraction is untouched. Whether Tesseract can
read a particular font is Tesseract's problem, covered separately in
test_ocr.py.
"""

from __future__ import annotations

import io

import pytest

from app.ingestion import ocr
from app.ingestion.base import UnparsableDocumentError
from app.ingestion.pdf import PdfParser
from app.ingestion.powerpoint import PptxParser
from app.ingestion.word import DocxParser


def _png(size: tuple[int, int] = (600, 400), colour: str = "white") -> bytes:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", size, colour).save(buffer, format="PNG")
    return buffer.getvalue()


RECOVERED = "RECOVERED FROM IMAGE: quarterly revenue rose twelve percent"


@pytest.fixture
def ocr_reads(monkeypatch: pytest.MonkeyPatch):
    """OCR is available and returns a fixed string for every image."""

    def _install(text: str = RECOVERED) -> str:
        monkeypatch.setattr(ocr, "is_available", lambda: True)
        monkeypatch.setattr(ocr, "extract_text", lambda blob, context="": ocr.OcrResult(text=text))
        return text

    return _install


@pytest.fixture
def ocr_finds_nothing(monkeypatch: pytest.MonkeyPatch) -> None:
    """Available, but every image is blank. The common case in real files."""
    monkeypatch.setattr(ocr, "is_available", lambda: True)
    monkeypatch.setattr(
        ocr,
        "extract_text",
        lambda blob, context="": ocr.OcrResult(text="", rejected="too_short"),
    )


@pytest.fixture
def ocr_off(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(ocr, "is_available", lambda: False)


# ---------------------------------------------------------------------------
# Document builders
# ---------------------------------------------------------------------------


def _pdf_with_text_and_image() -> bytes:
    import pymupdf

    document = pymupdf.open()
    page = document.new_page()
    page.insert_text(
        (72, 100),
        "Deployment windows are Tuesday and Thursday between 10:00 and 16:00 UTC.",
        fontsize=11,
    )
    page.insert_image(pymupdf.Rect(72, 200, 372, 400), stream=_png())
    data = document.tobytes()
    document.close()
    return data


def _pdf_text_only() -> bytes:
    import pymupdf

    document = pymupdf.open()
    page = document.new_page()
    lines = [
        "Atlas Operations Handbook",
        "Releases ship on Tuesdays and Thursdays.",
        "Friday deploys are prohibited outright.",
        "Rollback means redeploying the previous image tag.",
    ]
    for offset, line in enumerate(lines):
        page.insert_text((72, 100 + offset * 20), line, fontsize=11)
    data = document.tobytes()
    document.close()
    return data


def _pdf_image_only() -> bytes:
    """A scanned page: one full-page image, no text layer at all."""
    import pymupdf

    document = pymupdf.open()
    page = document.new_page()
    page.insert_image(page.rect, stream=_png(size=(1200, 1600)))
    data = document.tobytes()
    document.close()
    return data


def _docx_with_text_and_image() -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph("Escalation pages the secondary after fifteen minutes.")
    document.add_picture(io.BytesIO(_png()))
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_with_text_and_image() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Capacity review"
    slide.shapes.add_picture(io.BytesIO(_png()), Inches(1), Inches(2), Inches(4), Inches(3))
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Existing behaviour must not change
# ---------------------------------------------------------------------------


def test_a_normal_text_pdf_is_unaffected_by_ocr(ocr_reads) -> None:
    """The headline regression risk. A page of prose with no images must
    behave exactly as it did before OCR existed."""
    ocr_reads()

    parsed = PdfParser().parse(_pdf_text_only(), filename="a.pdf")

    assert len(parsed.pages) == 1
    assert "Friday deploys are prohibited" in parsed.pages[0].text
    assert parsed.pages[0].metadata.get("content_type") is None


def test_text_extraction_is_identical_with_ocr_on_and_off(ocr_reads) -> None:
    off = PdfParser().parse(_pdf_text_only(), filename="a.pdf")
    ocr_reads()
    on = PdfParser().parse(_pdf_text_only(), filename="a.pdf")

    assert [page.text for page in on.pages] == [page.text for page in off.pages]


def test_ocr_never_runs_on_a_page_with_text_and_no_images(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Avoiding unnecessary OCR is a cost requirement, so it is asserted
    rather than assumed. Every page of every text document would otherwise
    pay for a subprocess launch."""
    calls: list[str] = []

    def _record(blob, context=""):
        calls.append(context)
        return ocr.OcrResult(text="x" * 60)

    monkeypatch.setattr(ocr, "is_available", lambda: True)
    monkeypatch.setattr(ocr, "extract_text", _record)

    PdfParser().parse(_pdf_text_only(), filename="a.pdf")

    assert calls == []


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def test_pdf_keeps_both_its_text_and_its_image_text(ocr_reads) -> None:
    recovered = ocr_reads()

    parsed = PdfParser().parse(_pdf_with_text_and_image(), filename="mixed.pdf")

    texts = [page.text for page in parsed.pages]
    assert any("Deployment windows are Tuesday" in text for text in texts)
    assert any(recovered in text for text in texts)


def test_pdf_ocr_text_is_marked_and_keeps_its_page_number(ocr_reads) -> None:
    ocr_reads()

    parsed = PdfParser().parse(_pdf_with_text_and_image(), filename="mixed.pdf")

    ocr_pages = [p for p in parsed.pages if p.metadata.get("content_type") == "ocr"]
    assert ocr_pages, "no page was marked as OCR-derived"
    assert all(page.number == 1 for page in ocr_pages)


def test_a_scanned_pdf_becomes_readable(ocr_reads) -> None:
    """The case that previously failed outright with "no text could be
    extracted"."""
    recovered = ocr_reads()

    parsed = PdfParser().parse(_pdf_image_only(), filename="scan.pdf")

    assert parsed.pages
    assert recovered in parsed.pages[0].text
    assert parsed.pages[0].metadata["content_type"] == "ocr"


def test_a_scanned_pdf_still_fails_cleanly_without_ocr(ocr_off) -> None:
    with pytest.raises(UnparsableDocumentError, match="No text could be extracted"):
        PdfParser().parse(_pdf_image_only(), filename="scan.pdf")


def test_an_image_with_no_readable_text_adds_nothing(ocr_finds_nothing) -> None:
    """A decorative photograph must not become an indexable page."""
    parsed = PdfParser().parse(_pdf_with_text_and_image(), filename="mixed.pdf")

    assert len(parsed.pages) == 1
    assert all(page.metadata.get("content_type") != "ocr" for page in parsed.pages)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def test_docx_keeps_both_its_text_and_its_image_text(ocr_reads) -> None:
    recovered = ocr_reads()

    parsed = DocxParser().parse(_docx_with_text_and_image(), filename="report.docx")

    texts = [page.text for page in parsed.pages]
    assert any("Escalation pages the secondary" in text for text in texts)
    assert any(recovered in text for text in texts)


def test_docx_ocr_text_carries_location_metadata(ocr_reads) -> None:
    ocr_reads()

    parsed = DocxParser().parse(_docx_with_text_and_image(), filename="report.docx")

    ocr_pages = [p for p in parsed.pages if p.metadata.get("content_type") == "ocr"]
    assert ocr_pages
    assert all(page.number >= 1 for page in ocr_pages)


def test_docx_without_ocr_is_unchanged(ocr_off) -> None:
    parsed = DocxParser().parse(_docx_with_text_and_image(), filename="report.docx")

    assert len(parsed.pages) == 1
    assert "Escalation pages the secondary" in parsed.pages[0].text


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def test_pptx_keeps_both_its_text_and_its_image_text(ocr_reads) -> None:
    recovered = ocr_reads()

    parsed = PptxParser().parse(_pptx_with_text_and_image(), filename="deck.pptx")

    texts = [page.text for page in parsed.pages]
    assert any("Capacity review" in text for text in texts)
    assert any(recovered in text for text in texts)


def test_pptx_ocr_text_keeps_its_slide_number(ocr_reads) -> None:
    """The slide number is the citation target, so it must survive OCR."""
    ocr_reads()

    parsed = PptxParser().parse(_pptx_with_text_and_image(), filename="deck.pptx")

    ocr_pages = [p for p in parsed.pages if p.metadata.get("content_type") == "ocr"]
    assert ocr_pages
    assert all(page.number == 1 for page in ocr_pages)


def test_pptx_without_ocr_is_unchanged(ocr_off) -> None:
    parsed = PptxParser().parse(_pptx_with_text_and_image(), filename="deck.pptx")

    assert len(parsed.pages) == 1
    assert "Capacity review" in parsed.pages[0].text
