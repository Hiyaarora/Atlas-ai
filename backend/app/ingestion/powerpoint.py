"""PPTX parsing via python-pptx."""

import io

from pptx import Presentation

from app.core.config import settings
from app.core.logging import get_logger
from app.ingestion import ocr
from app.ingestion.archives import guard_ooxml
from app.ingestion.base import DocumentParser, ParsedDocument, ParsedPage, UnparsableDocumentError
from app.ingestion.tables import render_table

logger = get_logger(__name__)


class PptxParser(DocumentParser):
    content_types = ("application/vnd.openxmlformats-officedocument.presentationml.presentation",)
    extensions = (".pptx",)

    def parse(self, data: bytes, *, filename: str) -> ParsedDocument:
        guard_ooxml(data, filename=filename)

        try:
            presentation = Presentation(io.BytesIO(data))
        except Exception as exc:  # noqa: BLE001
            raise UnparsableDocumentError(
                f"Could not read {filename} as a PowerPoint file. "
                "If it is an older .ppt file, save it as .pptx first."
            ) from exc

        pages: list[ParsedPage] = []
        ocr_budget = settings.ocr_max_images_per_document
        ocr_slides = 0

        # Slides map onto pages exactly — the cleanest citation target of any
        # format here. "Slide 12" is unambiguous in a way "page 12" of a Word
        # document never is.
        for number, slide in enumerate(presentation.slides, start=1):
            blocks: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # The title placeholder is the slide's heading; marking
                        # it keeps that structure through chunking.
                        is_title = shape == slide.shapes.title
                        blocks.append(f"# {text}" if is_title else text)

                elif shape.has_table:
                    rendered = render_table(
                        [[cell.text for cell in row.cells] for row in shape.table.rows]
                    )
                    if rendered:
                        blocks.append(rendered)

            # Speaker notes often carry the actual argument, while the slide
            # holds three bullet fragments. Dropping them loses the substance.
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    blocks.append(f"Speaker notes: {notes}")

            if blocks:
                pages.append(ParsedPage(number=number, text="\n\n".join(blocks)))

            # --- OCR ------------------------------------------------------
            # A slide's pictures often carry the substance: a screenshot of a
            # dashboard, a photographed whiteboard, a diagram whose labels are
            # the only text on it. The text frames above miss all of that.
            if ocr_budget > 0 and ocr.is_available():
                slide_text = "\n\n".join(blocks)
                recognised: list[str] = []

                for shape in slide.shapes:
                    if ocr_budget <= 0:
                        break
                    image = getattr(shape, "image", None)
                    if image is None:
                        continue
                    ocr_budget -= 1

                    try:
                        blob = image.blob
                    except Exception as exc:  # noqa: BLE001
                        logger.info(
                            "pptx_image_unreadable",
                            extra={"slide": number, "error": str(exc)},
                        )
                        continue

                    result = ocr.extract_text(blob, context=f"slide {number}")
                    if result.usable and not ocr.is_redundant(result.text, slide_text):
                        recognised.append(result.text)

                if recognised:
                    ocr_slides += 1
                    pages.append(
                        ParsedPage(
                            number=number,
                            text="\n\n".join(recognised),
                            # Slide number preserved, so a citation still
                            # points at the slide the image sits on.
                            metadata={"content_type": ocr.OCR_CONTENT_TYPE},
                        )
                    )

        if not pages:
            raise UnparsableDocumentError(
                f"No text could be extracted from {filename}. "
                + (
                    "The slides appear to be images with no recognisable text."
                    if ocr.is_available()
                    else "If the slides are images, OCR is required and is not available here."
                )
            )

        logger.info(
            "pptx_parsed",
            extra={"source_name": filename, "slides": len(pages), "ocr_slides": ocr_slides},
        )

        core = presentation.core_properties
        return ParsedDocument(
            pages=pages,
            metadata={
                "title": core.title or None,
                "author": core.author or None,
                "source_page_count": len(presentation.slides._sldIdLst),
            },
        )
